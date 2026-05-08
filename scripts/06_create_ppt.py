"""
UFC 328 Prediction — Script 06: Generate PowerPoint Presentation
================================================================
Writes: outputs/UFC328_Prediction.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "outputs")
OUTPUTS    = os.path.join(os.path.dirname(__file__), "..", "outputs")

# Colour palette
NAVY    = RGBColor(0x1A, 0x1A, 0x2E)
BLUE    = RGBColor(0x16, 0x21, 0x3E)
ACCENT  = RGBColor(0xE9, 0x4F, 0x37)   # red-orange
GOLD    = RGBColor(0xF5, 0xA6, 0x23)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
LGREY   = RGBColor(0xCC, 0xCC, 0xCC)


def set_slide_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_image_safe(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        add_textbox(slide, f"[image: {os.path.basename(path)}]",
                    left, top, width, height, font_size=11, color=LGREY)


def add_divider(slide, top, color=ACCENT):
    from pptx.util import Pt as PT
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.4), Inches(top), Inches(9.2), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # completely blank


# ══════════════════════════════════════════════
# SLIDE 1 — Title / Hook
# ══════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1, NAVY)

add_textbox(s1, "CAN A MACHINE LEARNING MODEL", 1, 0.6, 11, 1.0,
            font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, "PREDICT A UFC CHAMPIONSHIP FIGHT?", 1, 1.4, 11, 1.0,
            font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_divider(s1, 2.5)

add_textbox(s1, "Khamzat Chimaev  vs  Sean Strickland", 1, 2.7, 11, 0.7,
            font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_textbox(s1, "UFC 328  |  Middleweight Title  |  May 9, 2026  |  Newark NJ", 1, 3.3, 11, 0.5,
            font_size=14, color=LGREY, align=PP_ALIGN.CENTER)

add_divider(s1, 4.0)

stats = [
    ("8,337", "UFC fights\nin training data"),
    ("527",   "MW fights\nused to train"),
    ("13",    "differential\nfeatures"),
    ("0.712", "best model\nROC-AUC"),
]
for i, (num, lbl) in enumerate(stats):
    x = 1.2 + i * 2.8
    add_textbox(s1, num, x, 4.2, 2.2, 0.7,
                font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s1, lbl, x, 4.9, 2.2, 0.7,
                font_size=11, color=LGREY, align=PP_ALIGN.CENTER)

add_textbox(s1, "XGBoost  |  Real UFC Data  |  Full ML Pipeline", 1, 6.5, 11, 0.6,
            font_size=13, italic=True, color=LGREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# SLIDE 2 — The Two Fighters
# ══════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2, NAVY)

add_textbox(s2, "THE FIGHTERS", 0.4, 0.2, 12, 0.6,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_divider(s2, 0.95)

# Left — Chimaev
add_textbox(s2, "KHAMZAT CHIMAEV", 0.5, 1.1, 5.5, 0.6,
            font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s2, '"Borz"  |  Champion', 0.5, 1.65, 5.5, 0.4,
            font_size=13, color=GOLD, align=PP_ALIGN.CENTER)

chimaev_stats = [
    ("Record",        "15 - 0  (UNDEFEATED)"),
    ("Win Streak",    "12 consecutive wins"),
    ("Finish Rate",   "80%  (12 of 15 wins)"),
    ("TD Avg/15min",  "5.29  (elite grappler)"),
    ("Sub Avg/15min", "1.80  (major sub threat)"),
    ("SLpM",          "5.24"),
    ("Age",           "32.3 years"),
]
for i, (k, v) in enumerate(chimaev_stats):
    y = 2.2 + i * 0.55
    add_textbox(s2, k, 0.5, y, 2.4, 0.45, font_size=12, color=LGREY)
    add_textbox(s2, v, 3.0, y, 3.0, 0.45, font_size=12, bold=True, color=WHITE)

# Divider between fighters
div = s2.shapes.add_shape(1, Inches(6.5), Inches(1.0), Inches(0.04), Inches(5.8))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT
div.line.fill.background()

# Right — Strickland
add_textbox(s2, "SEAN STRICKLAND", 6.8, 1.1, 5.5, 0.6,
            font_size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(s2, '"Tarzan"  |  Former Champion  |  Challenger', 6.8, 1.65, 5.5, 0.4,
            font_size=13, color=GOLD, align=PP_ALIGN.CENTER)

strickland_stats = [
    ("Record",        "30 - 7"),
    ("Win Streak",    "1 win"),
    ("Finish Rate",   "35%"),
    ("TD Avg/15min",  "0.43  (striker, not grappler)"),
    ("Sub Avg/15min", "0.10"),
    ("SLpM",          "7.25  (high volume striker)"),
    ("Age",           "35.2 years"),
]
for i, (k, v) in enumerate(strickland_stats):
    y = 2.2 + i * 0.55
    add_textbox(s2, k, 6.8, y, 2.4, 0.45, font_size=12, color=LGREY)
    add_textbox(s2, v, 9.3, y, 3.5, 0.45, font_size=12, bold=True, color=WHITE)


# ══════════════════════════════════════════════
# SLIDE 3 — Pipeline + Feature Engineering
# ══════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3, NAVY)

add_textbox(s3, "THE ML PIPELINE", 0.4, 0.2, 12, 0.6,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_divider(s3, 0.95)

pipeline_path = os.path.join(OUTPUTS, "viz_05_pipeline.png")
add_image_safe(s3, pipeline_path, 0.4, 1.1, 12.5, 4.2)

add_textbox(s3, "Key design decision: win_streak and finish_rate computed chronologically — no data leakage",
            0.4, 5.5, 12.5, 0.6, font_size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# SLIDE 4 — What the Model Sees (Differentials)
# ══════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4, NAVY)

add_textbox(s4, "WHAT THE MODEL SEES", 0.4, 0.2, 12, 0.6,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_divider(s4, 0.95)

diff_path = os.path.join(OUTPUTS, "viz_04_differentials.png")
add_image_safe(s4, diff_path, 0.3, 1.1, 7.5, 5.8)

add_textbox(s4, "Why Differential Features?", 8.2, 1.2, 4.7, 0.5,
            font_size=16, bold=True, color=GOLD)
add_textbox(s4,
    "Raw stat:\nChimaev TDs = 5.29/15min\n\nMeaningless alone.\n\n"
    "Differential:\nChimaev TDs - Strickland TDs\n= 5.29 - 0.43 = +4.86\n\n"
    "NOW the model knows:\nChimaev has a +4.86\ngrappling edge.\n\n"
    "13 features.\nEach = a competitive gap,\nnot an absolute value.",
    8.2, 1.8, 4.7, 5.0, font_size=12, color=WHITE)


# ══════════════════════════════════════════════
# SLIDE 5 — Model Performance
# ══════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5, NAVY)

add_textbox(s5, "MODEL PERFORMANCE", 0.4, 0.2, 12, 0.6,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_divider(s5, 0.95)

model_path = os.path.join(OUTPUTS, "viz_01_model_comparison.png")
add_image_safe(s5, model_path, 0.3, 1.1, 7.5, 5.6)

add_textbox(s5, "Key Insight", 8.2, 1.2, 4.7, 0.5,
            font_size=16, bold=True, color=GOLD)
add_textbox(s5,
    "Logistic Regression (0.712)\noutperformed XGBoost (0.671).\n\n"
    "This means the relationship\nbetween differential features\n"
    "and fight outcomes is largely\nLINEAR.\n\n"
    "XGBoost's extra complexity\nwasn't rewarded on\n527 training samples.\n\n"
    "ROC-AUC of 0.71 is competitive\nwith published academic\nUFC prediction models\n(typical range: 0.65-0.75).",
    8.2, 1.8, 4.7, 5.0, font_size=12, color=WHITE)


# ══════════════════════════════════════════════
# SLIDE 6 — Prediction + Takeaways
# ══════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6, NAVY)

add_textbox(s6, "THE PREDICTION", 0.4, 0.2, 12, 0.6,
            font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_divider(s6, 0.95)

pred_path = os.path.join(OUTPUTS, "viz_02_prediction_result.png")
add_image_safe(s6, pred_path, 0.3, 1.1, 5.8, 4.4)

# Benchmark table
add_textbox(s6, "Benchmark Comparison", 6.5, 1.1, 6.4, 0.5,
            font_size=16, bold=True, color=GOLD)
benchmarks = [
    ("Source",            "Chimaev",  "Strickland"),
    ("This Model (XGB)",  "59.0%",    "41.0%"),
    ("Elo Model",         "79.0%",    "21.0%"),
    ("Betting Odds",      "~85%",     "~15%"),
]
for i, (src, c, s) in enumerate(benchmarks):
    y = 1.75 + i * 0.6
    col = GOLD if i == 0 else (GREEN if i == 1 else LGREY)
    add_textbox(s6, src, 6.5, y, 3.5, 0.5, font_size=12, bold=(i==0), color=col)
    add_textbox(s6, c,   10.0, y, 1.5, 0.5, font_size=12, bold=(i==0), color=col, align=PP_ALIGN.CENTER)
    add_textbox(s6, s,   11.5, y, 1.5, 0.5, font_size=12, bold=(i==0), color=col, align=PP_ALIGN.CENTER)

add_divider(s6, 4.3)

add_textbox(s6, "Key Takeaways", 0.4, 4.5, 12.5, 0.5,
            font_size=16, bold=True, color=GOLD)
takeaways = [
    "Real data always beats synthetic data.",
    "Simple models can outperform complex ones on small datasets.",
    "Differential features > raw features for head-to-head prediction.",
    "No leakage: win_streak computed chronologically across 8,337 fights.",
]
for i, t in enumerate(takeaways):
    add_textbox(s6, f"  {t}", 0.4, 5.05 + i * 0.47, 12.5, 0.42,
                font_size=13, color=WHITE)

add_textbox(s6, "Full code on GitHub  |  #MachineLearning #DataScience #UFC #XGBoost",
            0.4, 7.05, 12.5, 0.35, font_size=11, italic=True, color=LGREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
out_path = os.path.join(OUTPUT_DIR, "UFC328_Prediction.pptx")
prs.save(out_path)
print(f"Saved -> {out_path}")
